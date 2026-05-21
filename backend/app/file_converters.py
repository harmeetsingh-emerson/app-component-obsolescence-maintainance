"""
File Converters — uniform BOM extraction for all supported file types.

Supported formats:
  PDF          → handled by bom_parser_v2 / ocr_processor (existing)
  TXT          → handled by ocr_processor (existing)
  Excel        → .xlsx / .xls  (openpyxl)
  CSV          → built-in csv module
  Word         → .docx  (python-docx)
  PowerPoint   → .pptx  (python-pptx)
  Images       → .png .jpg .jpeg .bmp .tiff .tif .gif .webp  (PaddleOCR)

All tabular formats (Excel, CSV, Word, PPT) are converted to the same
list-of-lists structure consumed by detect_bom_structure / parse_bom_row
from bom_parser_v2, so BOM recognition logic is reused without duplication.

Each public function returns a list of part dicts matching the shape
produced by parse_bom_document(), or raises an exception on hard failure.
"""

import csv
import json
import os
import re
import requests
from typing import List, Dict, Optional, Tuple

_OLLAMA_BASE = os.environ.get("OLLAMA_BASE_URL", "http://localhost:11434").rstrip("/")
_OLLAMA_CHAT_URL = f"{_OLLAMA_BASE}/api/chat"
_ANALYSIS_MODEL = "llama3.2:3b"


# ---------------------------------------------------------------------------
# LLM-assisted file structure analysis
# ---------------------------------------------------------------------------

_ANALYSIS_SYSTEM_PROMPT = """\
You are a file structure analyst for an electronics component obsolescence system.
You will be given column headers and a few sample rows from a spreadsheet or document table.

Your job is to:
1. Understand what the file/sheet is about.
2. Identify which column index contains each of these semantic fields (use null if absent):
   - part_number  : internal part number / BOM number / item number
   - manufacturer : manufacturer / vendor / brand name / mfr / mfr name / make / supplier
   - mpn          : manufacturer part number / MPN / catalog number / xref
   - description  : component description / subclass / type / value
   - quantity     : quantity / qty / count
   - designators  : reference designators / placement / location
3. Decide if this is a real BOM/component data table (is_bom_data: true) or a
   metadata/summary/info sheet (is_bom_data: false).

Return ONLY a valid JSON object — no explanation, no markdown fences:
{
  "is_bom_data": <true|false>,
  "file_topic": "<one sentence describing what this sheet contains>",
  "header_row_idx": <0-based index of the header row>,
  "columns": {
    "part_number":  <col index or null>,
    "manufacturer": <col index or null>,
    "mpn":          <col index or null>,
    "description":  <col index or null>,
    "quantity":     <col index or null>,
    "designators":  <col index or null>
  },
  "important_columns": [<indices of columns that matter most for queries>]
}"""


def _llm_analyze_table(headers: List[str], sample_rows: List[List[str]]) -> Dict:
    """
    Send headers + sample rows to the LLM and get back a structured analysis:
    which columns map to which semantic fields, and whether this is BOM data.
    Falls back to a heuristic result if the LLM is unavailable.
    """
    # Build a compact text representation of the table
    header_line = " | ".join(f"[{i}] {h}" for i, h in enumerate(headers) if h.strip())
    sample_lines = []
    for r in sample_rows[:5]:
        row_vals = " | ".join(str(r[i]).strip() if i < len(r) else "" for i in range(len(headers)))
        sample_lines.append(row_vals)

    user_msg = f"Headers:\n{header_line}\n\nSample rows:\n" + "\n".join(sample_lines)

    default = {
        "is_bom_data": False,
        "file_topic": "unknown",
        "header_row_idx": 0,
        "columns": {k: None for k in ("part_number","manufacturer","mpn","description","quantity","designators")},
        "important_columns": [],
    }

    try:
        payload = {
            "model": _ANALYSIS_MODEL,
            "messages": [
                {"role": "system", "content": _ANALYSIS_SYSTEM_PROMPT},
                {"role": "user",   "content": user_msg},
            ],
            "stream": False,
            "options": {"temperature": 0, "num_predict": 400},
        }
        resp = requests.post(_OLLAMA_CHAT_URL, json=payload, timeout=30)
        if resp.status_code != 200:
            print(f"[LLMAnalyzer] Ollama HTTP {resp.status_code} — using heuristic fallback")
            return default

        content = resp.json().get("message", {}).get("content", "").strip()
        json_match = re.search(r'\{.*\}', content, re.DOTALL)
        if not json_match:
            print(f"[LLMAnalyzer] No JSON in response — using heuristic fallback")
            return default

        result = {**default, **json.loads(json_match.group())}
        print(f"[LLMAnalyzer] topic={result['file_topic']!r}  is_bom={result['is_bom_data']}  cols={result['columns']}")
        return result

    except Exception as exc:
        print(f"[LLMAnalyzer] Error: {exc} — using heuristic fallback")
        return default


def _parse_table_intelligent(
    table: List[List[str]],
    sheet_label: str,
    sheet_idx: int = 1,
) -> List[Dict]:
    """
    LLM-guided table parser.

    Pipeline:
    1. Scan first 20 rows to find the best header row (most non-empty cells).
    2. Send headers + 5 sample rows to the LLM for structure analysis.
    3. Use the LLM's column mapping to read every data row.
    4. For BOM data: build proper part dicts with part_number / manufacturer / mpn.
       For non-BOM data: skip the sheet entirely.
    5. Each row is also stored as a full "Header: Value" text for semantic search.
    """
    if not table or len(table) < 2:
        return []

    # ── Step 1: Find the best header row ─────────────────────────────────
    best_idx, best_count = 0, 0
    for i, row in enumerate(table[:20]):
        non_empty = sum(1 for c in row if str(c).strip())
        if non_empty > best_count:
            best_count = non_empty
            best_idx = i

    headers = [str(c).strip() for c in table[best_idx]]
    data_rows = table[best_idx + 1:]

    if best_count < 2:
        print(f"[LLMAnalyzer] {sheet_label!r}: too few columns ({best_count}) — skipping")
        return []

    # ── Step 2: LLM analysis ─────────────────────────────────────────────
    non_blank_data = [r for r in data_rows if any(str(c).strip() for c in r)]
    analysis = _llm_analyze_table(headers, non_blank_data[:5])

    if not analysis.get("is_bom_data"):
        print(f"[LLMAnalyzer] {sheet_label!r}: LLM says not BOM data ({analysis.get('file_topic')}) — skipping")
        return []

    col_map = analysis.get("columns", {})
    # All non-empty column indices matter for the full row text
    all_col_indices = [i for i, h in enumerate(headers) if h]

    SKIP_VALUES = {"", "none", "n/a", "-", "nan"}

    # ── Step 2.5: Detect numbered manufacturer columns in headers ─────────
    # e.g. "Manufacturer 1", "Manufacturer Part Number 1", "MANUFACTURER 2", etc.
    # When present, we use them directly instead of the LLM's single-column mapping
    # (the LLM can't reliably pick ONE column when there are 6 mfr-related columns).
    import re as _re

    _MFR_COL_RE = _re.compile(
        r'^(?:manufacturer|mfr)\s*(\d+)$', _re.IGNORECASE
    )
    _MPN_COL_RE = _re.compile(
        r'^(?:manufacturer\s*part\s*(?:number|no|#?)|mfr\s*part\s*(?:number|no|#?)|mpn)\s*(\d+)$',
        _re.IGNORECASE
    )

    numbered_mfr_cols: Dict[int, Dict[str, int]] = {}  # {n: {"mfr": col_idx, "mpn": col_idx}}
    for i, hdr in enumerate(headers):
        m = _MFR_COL_RE.match(hdr.strip())
        if m:
            numbered_mfr_cols.setdefault(int(m.group(1)), {})["mfr"] = i
            continue
        m = _MPN_COL_RE.match(hdr.strip())
        if m:
            numbered_mfr_cols.setdefault(int(m.group(1)), {})["mpn"] = i

    # All column indices used by numbered manufacturer columns — excluded from extra_fields
    numbered_mfr_col_indices: set = set()
    for col_info in numbered_mfr_cols.values():
        numbered_mfr_col_indices.update(col_info.values())

    # Structured column indices (from LLM map) — also excluded from extra_fields.
    # When numbered mfr columns exist, drop the LLM's manufacturer/mpn mappings
    # so they don't conflict with the numbered ones.
    llm_structured_keys = ["part_number", "description", "quantity", "designators"]
    if not numbered_mfr_cols:
        llm_structured_keys += ["manufacturer", "mpn"]
    structured_indices = {
        col_map.get(k) for k in llm_structured_keys if col_map.get(k) is not None
    } | numbered_mfr_col_indices

    # ── Step 3: Parse every data row ─────────────────────────────────────
    parts: List[Dict] = []
    for row_idx, row in enumerate(data_rows, start=1):
        cells = [str(row[i]).strip() if i < len(row) else "" for i in range(len(headers))]
        meaningful = [c for c in cells if c.lower() not in SKIP_VALUES]
        if len(meaningful) < 2:
            continue

        def get(col_key):
            idx = col_map.get(col_key)
            if idx is None or idx >= len(cells):
                return ""
            v = cells[idx]
            return v if v.lower() not in SKIP_VALUES else ""

        # Full row text: "Header: Value | ..." for all non-empty cols
        pairs = []
        for i in all_col_indices:
            val = cells[i] if i < len(cells) else ""
            if val and val.lower() not in SKIP_VALUES:
                pairs.append(f"{headers[i]}: {val}")
        row_text = " | ".join(pairs)

        # Part number — required
        part_number = get("part_number")
        if not part_number:
            part_number = next((c for c in cells if c.lower() not in SKIP_VALUES),
                               f"{sheet_label}-R{row_idx}")

        description = get("description") or row_text
        quantity    = get("quantity")
        designators = get("designators")
        designators  = get("designators")

        # ── Build manufacturers list ───────────────────────────────────────
        manufacturers_list = []

        if numbered_mfr_cols:
            # Numbered columns exist — build list deterministically from them.
            # Completely ignore the LLM's manufacturer/mpn mapping here.
            for n in sorted(numbered_mfr_cols):
                col_info = numbered_mfr_cols[n]
                mfr_val = (cells[col_info["mfr"]] if "mfr" in col_info and col_info["mfr"] < len(cells) else "").strip()
                mpn_val = (cells[col_info["mpn"]] if "mpn" in col_info and col_info["mpn"] < len(cells) else "").strip()
                if mfr_val.lower() in SKIP_VALUES:
                    mfr_val = ""
                if mpn_val.lower() in SKIP_VALUES:
                    mpn_val = ""
                if mfr_val or mpn_val:
                    manufacturers_list.append({
                        "manufacturer": mfr_val or "UNKNOWN",
                        "mpn": mpn_val or "",
                        "preference": n,
                        "confidence": 0.8,
                    })
        else:
            # No numbered pattern — use LLM's single manufacturer/mpn mapping.
            manufacturer = get("manufacturer")
            mpn          = get("mpn")
            if manufacturer or mpn:
                manufacturers_list.append({
                    "manufacturer": manufacturer or "UNKNOWN",
                    "mpn": mpn or "",
                    "preference": 1,
                    "confidence": 0.7,
                })

        # ── extra_fields: every column not already in a structured field ──
        # Store once with original-case key only — no lowercase duplicates.
        extra_fields: Dict[str, str] = {}
        for i, hdr in enumerate(headers):
            if not hdr or i in structured_indices:
                continue
            val = cells[i] if i < len(cells) else ""
            if val and val.lower() not in SKIP_VALUES:
                extra_fields[hdr] = val

        parts.append({
            "part_number":   part_number,
            "description":   description,
            "raw_row_text":  row_text,
            "manufacturers": manufacturers_list,
            "quantity":      quantity,
            "designators":   designators,
            "extra_fields":  extra_fields,
            "confidence":    0.8,
            "page_number":   sheet_idx,
            "source_type":   "intelligent_row",
            "sheet":         sheet_label,
            "file_topic":    analysis.get("file_topic", ""),
        })

    print(f"[LLMAnalyzer] {sheet_label!r}: indexed {len(parts)} rows")
    return parts


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _save_manufacturer_db() -> None:
    """Save the in-memory manufacturer database (mirrors PDF parser behaviour)."""
    try:
        from app.bom_parser_v2 import save_manufacturers, get_manufacturer_stats
        stats = get_manufacturer_stats()
        if stats.get('learned', 0) > 0:
            print(f"[Manufacturer DB] Saving {stats['learned']} learned manufacturer(s)...")
            if save_manufacturers():
                print("[Manufacturer DB] Saved OK")
    except Exception as exc:
        print(f"[Manufacturer DB] Save failed: {exc}")


def _parse_table_to_parts(
    table: List[List[str]],
    source_label: str,
    page_num: int = 1,
    last_mapping=None,
) -> tuple:
    """
    Run the BOM detection + row parsing pipeline on a raw list-of-lists table.
    Mirrors the PDF parser pipeline from parse_bom_document:
      - detect_bom_structure (with lower threshold for tabular sources)
      - _infer_columns_from_content for missing mpn/manufacturer
      - last_mapping fallback (same as PDF cross-page persistence)
      - parse_bom_row for each data row

    Returns (parts: List[Dict], detection_used)
    detection_used is the BOMTableDetection that succeeded, or None if nothing worked.
    """
    from app.bom_parser_v2 import (
        detect_bom_structure, parse_bom_row, _infer_columns_from_content
    )

    if not table:
        return [], last_mapping

    detection = detect_bom_structure(table, page_num=page_num)

    # For non-PDF tabular sources the header is always present — accept
    # borderline confidence (≥0.2 instead of the PDF default 0.4).
    if not detection.is_bom and detection.confidence >= 0.2 and detection.header_row_idx >= 0:
        print(f"[FileConverter] {source_label}: low confidence ({detection.confidence:.2f}) "
              f"but header found — treating as BOM anyway")
        detection.is_bom = True

    if not detection.is_bom:
        print(f"[FileConverter] {source_label}: detection failed "
              f"(confidence={detection.confidence:.2f}, header_row={detection.header_row_idx})")

        # ── Mirror PDF parser: reuse mapping from a previous sheet/table ──
        if last_mapping:
            print(f"[FileConverter] {source_label}: reusing column map from previous table")
            detection = last_mapping
            detection.header_row_idx = -1   # whole table is data rows
        else:
            # Last resort: force row 0 as header + content inference
            if len(table) >= 2:
                print(f"[FileConverter] {source_label}: forced header=row0 + content inference")
                n_cols = max((len(r) for r in table if r), default=0)
                col_map = _infer_columns_from_content(table, 0, {}, n_cols)
                if 'mpn' in col_map or 'manufacturer' in col_map:
                    detection.is_bom = True
                    detection.header_row_idx = 0
                    detection.column_map = col_map
                    print(f"[FileConverter] {source_label}: forced column map: {col_map}")
                else:
                    return [], last_mapping
            else:
                return [], last_mapping

    # Run content inference to fill in any missing mpn / manufacturer columns
    n_cols = max((len(r) for r in table if r), default=0)
    detection.column_map = _infer_columns_from_content(
        table, detection.header_row_idx, detection.column_map, n_cols
    )

    print(f"[FileConverter] {source_label}: column map = {detection.column_map}")

    data_rows = table[detection.header_row_idx + 1:] if detection.header_row_idx >= 0 else table
    parts: List[Dict] = []

    for row in data_rows:
        part = parse_bom_row(row, detection.column_map, page_num)
        if part:
            part_dict = {
                "part_number": part.part_number,
                "description": part.description,
                "manufacturers": part.manufacturers,
                "quantity": part.quantity,
                "designators": part.designators,
                "confidence": part.confidence,
                "page_number": page_num,
            }
            if part.validation_flags:
                part_dict["validation_flags"] = part.validation_flags
            parts.append(part_dict)

    print(f"[FileConverter] {source_label} — page/sheet {page_num}: {len(parts)} parts")
    return parts, detection


# ---------------------------------------------------------------------------
# Raw-row fallback — used when structured BOM parsing finds nothing
# ---------------------------------------------------------------------------

# Keywords that indicate a row is likely a BOM column header row
_BOM_HEADER_KEYWORDS = {
    'part', 'pn', 'p/n', 'mpn', 'mfr', 'mfgr', 'manufacturer',
    'description', 'desc', 'qty', 'quantity', 'designator', 'ref',
    'reference', 'item', 'bom', 'component', 'value', 'footprint',
    'package', 'comment', 'subclass', 'class', 'category', 'level',
    'find', 'xref', 'cell', 'specification', 'spec', 'type',
}

# Keywords that identify a sheet as a metadata / summary sheet (skip it)
_META_SHEET_KEYWORDS = {
    'project full path', 'project filename', 'variant name',
    'data-source', 'report time', 'report date', 'output name',
    'output type', 'output generator', 'schematic info',
}


def _find_bom_header_row(table: List[List[str]], max_scan: int = 20):
    """
    Scan the first `max_scan` rows and return (row_index, score) for the row
    that looks most like a BOM column-header row.
    Returns (-1, 0) if no suitable row is found.
    """
    best_idx = -1
    best_score = 0

    for row_idx, row in enumerate(table[:max_scan]):
        cells = [str(c).strip() for c in row if str(c).strip()]
        if len(cells) < 3:          # need at least 3 non-empty columns
            continue
        score = 0
        for c in cells:
            cl = c.lower()
            # direct keyword match
            if cl in _BOM_HEADER_KEYWORDS:
                score += 2
            # substring match (handles "Part Number", "Mfr Part", etc.)
            elif any(kw in cl for kw in _BOM_HEADER_KEYWORDS):
                score += 1
        if score > best_score:
            best_score = score
            best_idx = row_idx

    return best_idx, best_score


def _is_metadata_sheet(table: List[List[str]], max_scan: int = 15) -> bool:
    """
    Return True if the sheet looks like a metadata/summary sheet rather than
    a BOM data table.  Heuristics:
    - Fewer than 3 non-empty columns on the best candidate header row, OR
    - First column contains known metadata labels.
    """
    # Check first column for metadata labels
    first_col_hits = 0
    for row in table[:max_scan]:
        if not row:
            continue
        first_cell = str(row[0]).strip().lower()
        if first_cell in _META_SHEET_KEYWORDS or any(kw in first_cell for kw in _META_SHEET_KEYWORDS):
            first_col_hits += 1
    if first_col_hits >= 3:
        return True

    # Check max column width across first few rows
    max_cols = max((sum(1 for c in row if str(c).strip()) for row in table[:max_scan] if row), default=0)
    if max_cols < 3:
        return True

    return False


def _table_to_raw_row_parts(
    table: List[List[str]],
    sheet_label: str,
    sheet_idx: int = 1,
) -> List[Dict]:
    """
    Convert a raw table into part dicts for FAISS indexing.

    - Automatically finds the real header row (not always row 0).
    - Skips metadata/summary sheets.
    - Each non-blank data row becomes one FAISS record with full
      "Header: Value" text as the searchable description.
    """
    if not table or len(table) < 2:
        return []

    # Skip sheets that are clearly metadata/summary sheets
    if _is_metadata_sheet(table):
        print(f"[FileConverter] Skipping metadata sheet: {sheet_label!r}")
        return []

    # Find the best header row
    header_row_idx, score = _find_bom_header_row(table)
    if header_row_idx < 0 or score < 2:
        print(f"[FileConverter] No BOM header found in {sheet_label!r} (score={score}) — skipping")
        return []

    headers = [str(c).strip() for c in table[header_row_idx]]
    non_empty_cols = [i for i, h in enumerate(headers) if h]

    print(f"[FileConverter] {sheet_label!r}: header at row {header_row_idx} "
          f"(score={score}), {len(non_empty_cols)} columns: {[headers[i] for i in non_empty_cols]}")

    # Identify the most likely "part number" column
    PN_HINTS = {
        'part number', 'part no', 'p/n', 'pn', 'item', 'item no',
        'bom no', 'find no', 'component no', 'material no', 'part id',
        'designator', 'ref', 'reference',
    }
    pn_col = None
    for i in non_empty_cols:
        h = headers[i].lower()
        if h in PN_HINTS or any(hint in h for hint in PN_HINTS):
            pn_col = i
            break

    SKIP_VALUES = {"", "none", "n/a", "-", "nan", "true", "false"}

    parts: List[Dict] = []
    for row_idx, row in enumerate(table[header_row_idx + 1:], start=1):
        cells = [str(row[i]).strip() if i < len(row) else "" for i in range(len(headers))]

        # Skip blank / trivial rows
        meaningful = [c for c in cells if c.lower() not in SKIP_VALUES]
        if len(meaningful) < 2:
            continue

        # Build "Header: Value" pairs
        pairs = []
        for i in non_empty_cols:
            val = cells[i] if i < len(cells) else ""
            if val and val.lower() not in SKIP_VALUES:
                pairs.append(f"{headers[i]}: {val}")

        if not pairs:
            continue

        row_text = " | ".join(pairs)

        # Best-effort part number
        if pn_col is not None and pn_col < len(cells) and cells[pn_col].lower() not in SKIP_VALUES:
            part_number = cells[pn_col]
        else:
            part_number = next(
                (c for c in cells if c and c.lower() not in SKIP_VALUES),
                f"{sheet_label}-R{row_idx}",
            )

        # Build extra_fields: header → value, original case only (no lowercase dupes)
        extra_fields: Dict[str, str] = {}
        for i in non_empty_cols:
            val = cells[i] if i < len(cells) else ""
            if val and val.lower() not in SKIP_VALUES:
                extra_fields[headers[i]] = val

        parts.append({
            "part_number": part_number,
            "description": row_text,
            "raw_row_text": row_text,
            "manufacturers": [],
            "quantity": "",
            "designators": "",
            "extra_fields": extra_fields,
            "confidence": 0.5,
            "page_number": sheet_idx,
            "source_type": "raw_row",
            "sheet": sheet_label,
        })

    print(f"[FileConverter] Raw-row fallback: {len(parts)} rows from {sheet_label!r}")
    return parts


def parse_excel_raw_rows(file_path: str) -> List[Dict]:
    """
    LLM-guided Excel reader.
    For each sheet: loads the table, asks the LLM what the sheet is about and
    which columns matter, then indexes every data row intelligently.
    Replaces the old keyword-based raw-row fallback.
    """
    import openpyxl

    print(f"[FileConverter] Intelligent Excel reading: {file_path}")
    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    all_parts: List[Dict] = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        table: List[List[str]] = []
        for row in ws.iter_rows(values_only=True):
            table.append([str(c) if c is not None else "" for c in row])

        if all(all(cell == "" for cell in row) for row in table):
            print(f"[FileConverter] Skipping blank sheet: {sheet_name!r}")
            continue

        parts = _parse_table_intelligent(table, sheet_name, sheet_idx)
        all_parts.extend(parts)

    print(f"[FileConverter] Intelligent total: {len(all_parts)} rows across {len(wb.sheetnames)} sheet(s)")
    return all_parts


# ---------------------------------------------------------------------------
# Excel  (.xlsx / .xls)
# ---------------------------------------------------------------------------

def parse_excel_bom(file_path: str) -> List[Dict]:
    """
    Excel parser — tries structured BOM detection first, then falls back to
    LLM-guided intelligent row reading per sheet.
    """
    import openpyxl
    from app.bom_parser_v2 import save_manufacturers, get_manufacturer_stats

    print(f"[FileConverter] Opening Excel: {file_path}")
    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    all_parts: List[Dict] = []
    last_mapping = None

    for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        print(f"[FileConverter] Processing sheet '{sheet_name}'")

        table: List[List[str]] = []
        for row in ws.iter_rows(values_only=True):
            table.append([str(c) if c is not None else "" for c in row])

        if all(all(cell == "" for cell in row) for row in table):
            continue

        # Try structured BOM detection first
        parts, last_mapping = _parse_table_to_parts(
            table, f"Excel sheet '{sheet_name}'", page_num=sheet_idx,
            last_mapping=last_mapping,
        )

        if not parts:
            # Structured parser found nothing — ask the LLM to figure it out
            print(f"[FileConverter] Structured parse empty for '{sheet_name}' — using LLM analysis")
            parts = _parse_table_intelligent(table, sheet_name, sheet_idx)

        all_parts.extend(parts)

    print(f"[FileConverter] Excel total: {len(all_parts)} parts across {len(wb.sheetnames)} sheet(s)")
    _save_manufacturer_db()
    return all_parts


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------

def parse_csv_bom(file_path: str) -> List[Dict]:
    """
    Parse BOM data from a CSV file.
    """
    print(f"[FileConverter] Opening CSV: {file_path}")
    table: List[List[str]] = []

    # Sniff the dialect so we handle both comma and semicolon separators
    with open(file_path, newline="", encoding="utf-8-sig") as fh:
        sample = fh.read(4096)
        fh.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            dialect = csv.excel  # fallback

        reader = csv.reader(fh, dialect)
        for row in reader:
            table.append([cell.strip() for cell in row])

    parts, _ = _parse_table_to_parts(table, "CSV", page_num=1)

    if not parts:
        print("[FileConverter] CSV structured parse empty — using LLM analysis")
        parts = _parse_table_intelligent(table, "CSV", sheet_idx=1)

    print(f"[FileConverter] CSV total: {len(parts)} parts")
    _save_manufacturer_db()
    return parts


# ---------------------------------------------------------------------------
# Word  (.docx)
# ---------------------------------------------------------------------------

def parse_word_bom(file_path: str) -> List[Dict]:
    """
    Parse BOM data from a Word document.
    Every table in the document is checked for BOM structure.
    """
    import docx  # python-docx
    from app.bom_parser_v2 import save_manufacturers, get_manufacturer_stats

    print(f"[FileConverter] Opening Word document: {file_path}")
    doc = docx.Document(file_path)
    all_parts: List[Dict] = []
    last_mapping = None

    for tbl_idx, table in enumerate(doc.tables, start=1):
        print(f"[FileConverter] Processing Word table {tbl_idx}")
        rows: List[List[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])

        if not rows:
            continue

        parts, last_mapping = _parse_table_to_parts(
            rows, f"Word table {tbl_idx}", page_num=tbl_idx,
            last_mapping=last_mapping,
        )
        if not parts:
            print(f"[FileConverter] Word table {tbl_idx} structured parse empty — using LLM analysis")
            parts = _parse_table_intelligent(rows, f"Word table {tbl_idx}", tbl_idx)
        all_parts.extend(parts)

    if not all_parts:
        # Fallback: collect all paragraph text and try OCR-style parsing
        print("[FileConverter] No BOM tables found in Word doc — trying paragraph text")
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if text:
            from app.ocr_processor import parse_ocr_bom_text
            all_parts = parse_ocr_bom_text(text)
            print(f"[FileConverter] Word paragraph fallback: {len(all_parts)} parts")

    print(f"[FileConverter] Word total: {len(all_parts)} parts")
    _save_manufacturer_db()
    return all_parts


# ---------------------------------------------------------------------------
# PowerPoint  (.pptx)
# ---------------------------------------------------------------------------

def parse_pptx_bom(file_path: str) -> List[Dict]:
    """
    Parse BOM data from a PowerPoint file.
    Tables in each slide are checked for BOM structure.
    Text from all shapes is collected as a fallback.
    """
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    print(f"[FileConverter] Opening PowerPoint: {file_path}")
    prs = Presentation(file_path)
    all_parts: List[Dict] = []
    all_text_lines: List[str] = []
    last_mapping = None

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"[FileConverter] Processing slide {slide_idx}")

        for shape in slide.shapes:
            # Collect text from every shape for fallback
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        all_text_lines.append(line)

            # Check for tables
            if shape.has_table:
                tbl = shape.table
                rows: List[List[str]] = []
                for row in tbl.rows:
                    rows.append([cell.text.strip() for cell in row.cells])

                parts, last_mapping = _parse_table_to_parts(
                    rows, f"Slide {slide_idx} table", page_num=slide_idx,
                    last_mapping=last_mapping,
                )
                all_parts.extend(parts)

    if not all_parts and all_text_lines:
        # Fallback: treat all collected text as OCR output
        print("[FileConverter] No BOM tables in PPT — trying text fallback")
        from app.ocr_processor import parse_ocr_bom_text
        all_parts = parse_ocr_bom_text("\n".join(all_text_lines))
        print(f"[FileConverter] PPT text fallback: {len(all_parts)} parts")

    print(f"[FileConverter] PowerPoint total: {len(all_parts)} parts")
    _save_manufacturer_db()
    return all_parts


# ---------------------------------------------------------------------------
# Images  (.png .jpg .jpeg .bmp .tiff .tif .gif .webp)
# ---------------------------------------------------------------------------

def parse_image_bom(file_path: str) -> List[Dict]:
    """
    Run PaddleOCR directly on an image file and parse the extracted text
    for BOM data.
    Returns a list of part dicts (may be empty for non-BOM images).
    """
    import numpy as np
    from PIL import Image
    from app.ocr_processor import _get_paddle_ocr_engine, parse_ocr_bom_text

    print(f"[FileConverter] Running OCR on image: {file_path}")
    engine = _get_paddle_ocr_engine()
    if engine is None:
        raise RuntimeError("PaddleOCR engine could not be initialised")

    # Load image via Pillow and convert to numpy array (RGB)
    img = Image.open(file_path).convert("RGB")
    img_array = np.array(img)

    raw = list(engine.predict(img_array))

    # Assemble OCR output into lines of text
    lines: List[str] = []
    for page_result in raw:
        if page_result is None:
            continue
        for item in page_result:
            if not item:
                continue
            try:
                # PaddleOCR v3: each item is a dict with 'rec_text' / 'rec_score'
                if isinstance(item, dict):
                    text = item.get("rec_text", "").strip()
                else:
                    # v2-style: [[bbox], (text, score)]
                    text = str(item[1][0]).strip() if len(item) > 1 else ""
                if text:
                    lines.append(text)
            except Exception:
                pass

    ocr_text = "\n".join(lines)
    print(f"[FileConverter] Image OCR: {len(ocr_text)} chars from {len(lines)} lines")

    if not ocr_text.strip():
        return []

    parts = parse_ocr_bom_text(ocr_text)
    print(f"[FileConverter] Image total: {len(parts)} parts")
    return parts


# ---------------------------------------------------------------------------
# Dispatch helper — called from the upload route
# ---------------------------------------------------------------------------

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".gif", ".webp"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}
SUPPORTED_EXTENSIONS = IMAGE_EXTENSIONS | EXCEL_EXTENSIONS | {".csv", ".docx", ".pptx", ".doc", ".ppt"}


def get_file_type(filename: str) -> Optional[str]:
    """
    Return a type token for the filename's extension, or None if unsupported.
    Tokens: 'excel', 'csv', 'word', 'pptx', 'image'
    PDF and TXT are handled by the existing routes and are not returned here.
    """
    ext = os.path.splitext(filename.lower())[1]
    if ext in EXCEL_EXTENSIONS:
        return "excel"
    if ext == ".csv":
        return "csv"
    if ext in {".docx", ".doc"}:
        return "word"
    if ext in {".pptx", ".ppt"}:
        return "pptx"
    if ext in IMAGE_EXTENSIONS:
        return "image"
    return None


def convert_and_parse(file_path: str, filename: str) -> List[Dict]:
    """
    Dispatch to the correct converter based on file extension.
    Returns a list of part dicts.
    Raises ValueError for unsupported types.
    """
    file_type = get_file_type(filename)
    if file_type == "excel":
        return parse_excel_bom(file_path)
    if file_type == "csv":
        return parse_csv_bom(file_path)
    if file_type == "word":
        return parse_word_bom(file_path)
    if file_type == "pptx":
        return parse_pptx_bom(file_path)
    if file_type == "image":
        return parse_image_bom(file_path)
    raise ValueError(f"Unsupported file type: {os.path.splitext(filename)[1]}")
